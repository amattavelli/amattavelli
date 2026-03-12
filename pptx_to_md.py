#!/usr/bin/env python3
"""
Converte un file PPTX in Markdown.
Uso: python3 pptx_to_md.py file.pptx [output.md]
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


def extract_text(shape):
    """Estrae il testo da una shape, ignorando shape vuote."""
    if not shape.has_text_frame:
        return []
    lines = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Determina il livello di rientro
        level = para.level
        if level == 0:
            lines.append(text)
        else:
            lines.append("  " * level + "- " + text)
    return lines


def pptx_to_markdown(pptx_path: Path) -> str:
    prs = Presentation(pptx_path)
    title = pptx_path.stem.replace("-", " ").replace("_", " ")
    lines = [f"# {title}", ""]

    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []

        # Separa titolo dal corpo
        title_shape = None
        body_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.shape_type == 13:  # immagine, skip
                    continue
                try:
                    is_title = title_shape is None and shape.placeholder_format is not None and shape.placeholder_format.idx == 0
                except ValueError:
                    is_title = False
                if is_title:
                    title_shape = shape
                else:
                    body_shapes.append(shape)

        # Titolo slide
        slide_title = ""
        if title_shape:
            slide_title = title_shape.text_frame.text.strip()
        if not slide_title:
            # Prova a prendere la prima shape con testo come titolo
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_title = shape.text_frame.text.strip().split("\n")[0]
                    break

        if slide_title:
            slide_lines.append(f"## Slide {i} — {slide_title}")
        else:
            slide_lines.append(f"## Slide {i}")

        # Corpo
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            texts = extract_text(shape)
            for t in texts:
                if t != slide_title:
                    slide_lines.append("- " + t if not t.startswith("-") and not t.startswith(" ") else t)

        if len(slide_lines) > 1:
            lines.extend(slide_lines)
            lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python3 pptx_to_md.py file.pptx [output.md]")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"File non trovato: {pptx_path}")
        sys.exit(1)

    md_content = pptx_to_markdown(pptx_path)

    if len(sys.argv) >= 3:
        out_path = Path(sys.argv[2])
    else:
        out_path = pptx_path.with_suffix(".md")

    out_path.write_text(md_content, encoding="utf-8")
    print(f"Creato: {out_path} ({len(md_content):,} caratteri)")
